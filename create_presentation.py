import sys
import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

template_path = 'C:/Users/Wuang/Desktop/Vibe/thuyetrinhslideluanvan.pptx'
output_path = 'C:/Users/Wuang/Desktop/Slide_ThuyetTrinh_LuanVan_WuangEvents.pptx'

prs = Presentation(template_path)

# Keep exactly 10 slides
while len(prs.slides) > 10:
    rId = prs.slides._sldIdLst[10].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[10]

print(f"Slide count truncated to: {len(prs.slides)}")

def clear_shapes_text(slide):
    """Clear text in textframes or prepare clean frames."""
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                p.text = ""

# COLOR PALETTE
COLOR_PRIMARY = RGBColor(15, 23, 42)      # Dark Navy #0F172A
COLOR_ACCENT = RGBColor(37, 99, 235)      # Blue #2563EB
COLOR_TEXT_MAIN = RGBColor(31, 41, 55)   # Charcoal #1F2937
COLOR_MUTED = RGBColor(100, 116, 139)    # Slate Gray #64748B

FONT_NAME = 'Segoe UI'

def add_header(slide, title_text, subtitle_text=None):
    """Utility to add clean title to slide."""
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_NAME
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.name = FONT_NAME
        p2.font.size = Pt(14)
        p2.font.color.rgb = COLOR_MUTED
        p2.space_before = Pt(4)

# ==========================================
# SLIDE 1: TRANG TIÊU ĐỀ
# ==========================================
slide1 = prs.slides[0]
clear_shapes_text(slide1)

tb1 = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(11.5), Inches(5.8))
tf1 = tb1.text_frame
tf1.word_wrap = True

p = tf1.paragraphs[0]
p.text = "TRƯỜNG ĐẠI HỌC CÔNG NGHỆ SÀI GÒN (STU)"
p.font.name = FONT_NAME
p.font.size = Pt(16)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

p = tf1.add_paragraph()
p.text = "KHOA CÔNG NGHỆ THÔNG TIN"
p.font.name = FONT_NAME
p.font.size = Pt(14)
p.font.bold = True
p.font.color.rgb = COLOR_MUTED
p.space_after = Pt(20)

p = tf1.add_paragraph()
p.text = "BÁO CÁO LUẬN VĂN TỐT NGHIỆP"
p.font.name = FONT_NAME
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY

p = tf1.add_paragraph()
p.text = "XÂY DỰNG WEBSITE QUẢN LÝ VÀ BÁN VÉ SỰ KIỆN TRỰC TUYẾN WUANGEVENTS"
p.font.name = FONT_NAME
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT
p.space_after = Pt(30)

p = tf1.add_paragraph()
p.text = "• Sinh viên thực hiện: Nguyễn Vinh Quang"
p.font.name = FONT_NAME
p.font.size = Pt(15)
p.font.bold = True
p.font.color.rgb = COLOR_TEXT_MAIN

p = tf1.add_paragraph()
p.text = "• Mã số sinh viên: DH52201310    |    Lớp: D22_TH08"
p.font.name = FONT_NAME
p.font.size = Pt(15)
p.font.color.rgb = COLOR_TEXT_MAIN

p = tf1.add_paragraph()
p.text = "• Chuyên ngành: Công nghệ Thông tin (Hệ thống Thông tin)"
p.font.name = FONT_NAME
p.font.size = Pt(15)
p.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 2: MỤC TIÊU ĐỀ TÀI
# ==========================================
slide2 = prs.slides[1]
clear_shapes_text(slide2)
add_header(slide2, "1. MỤC TIÊU ĐỀ TÀI", "Tóm tắt ngắn gọn các mục tiêu nghiên cứu và phát triển hệ thống")

tb2 = slide2.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf2 = tb2.text_frame
tf2.word_wrap = True

items2 = [
    ("Xây dựng Hệ thống Bán vé Trực tuyến toàn diện:", " Cung cấp giải pháp số hóa trọn gói cho thị trường sự kiện Việt Nam."),
    ("Tự động hóa Quy trình Bán vé & Giữ chỗ:", " Hỗ trợ giữ chỗ tạm thời 10 phút và phát hành Vé điện tử kèm mã QR Code độc bản."),
    ("Tích hợp 3 Cổng Thanh toán Điện tử:", " Kết nối trực tiếp API VNPay, MoMo và ZaloPay với cơ chế bảo mật chữ ký số HMAC-SHA256."),
    ("Bảo đảm An toàn Dữ liệu Đồng thời (Concurrency):", " Đạt chỉ số 0% Bán vượt kho (Overselling) và 0% Trùng ghế (Double Booking) khi tải cao."),
    ("Phân quyền Đa vai trò Chuyên biệt:", " Cung cấp 4 không gian quản trị độc lập cho Khách hàng, Ban tổ chức, Staff soát vé và Admin.")
]

for idx, (title, desc) in enumerate(items2):
    p = tf2.paragraphs[0] if idx == 0 else tf2.add_paragraph()
    p.space_after = Pt(12)
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT_NAME
    run2.font.size = Pt(15)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 3: NỘI DUNG THỰC HIỆN
# ==========================================
slide3 = prs.slides[2]
clear_shapes_text(slide3)
add_header(slide3, "2. NỘI DUNG THỰC HIỆN", "Quy trình khảo sát, thiết kế và phát triển dự án WuangEvents")

tb3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf3 = tb3.text_frame
tf3.word_wrap = True

items3 = [
    ("Phân tích Nghiệp vụ & Đặc tả Sử dụng:", " Khảo sát luồng nghiệp vụ bán vé thực tế, xây dựng 22 Use Cases và thiết kế CSDL chuẩn hóa (11 bảng)."),
    ("Phát triển Nền tảng Backend .NET 9:", " Sử dụng C# 13, ASP.NET Core MVC, Dapper Micro-ORM, mã hóa băm mật khẩu BCrypt và thư viện QRCoder."),
    ("Tích hợp API Cổng Thanh toán:", " Hiện thực luồng kết nối 3 cổng VNPay (SHA256 Checksum), MoMo (HMAC-SHA256) và ZaloPay (HMAC-SHA256)."),
    ("Thiết kế Giao diện Web Responsive:", " Xây dựng giao diện tương thích Máy tính & Điện thoại, tích hợp Sơ đồ chọn ghế trực quan (Seat Map)."),
    ("Kiểm thử & Đánh giá Hệ thống:", " Thực thi 27 Kịch bản Test Case (Chức năng, Bảo mật chữ ký số, Giữ chỗ 10 phút và Xử lý đồng thời).")
]

for idx, (title, desc) in enumerate(items3):
    p = tf3.paragraphs[0] if idx == 0 else tf3.add_paragraph()
    p.space_after = Pt(12)
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT_NAME
    run2.font.size = Pt(15)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 4: THÁCH THỨC VÀ GIẢI PHÁP
# ==========================================
slide4 = prs.slides[3]
clear_shapes_text(slide4)
add_header(slide4, "3. THÁCH THỨC VÀ GIẢI PHÁP KỸ THUẬT", "Các bài toán kỹ thuật phức tạp và giải pháp hiện thực trong ASP.NET Core")

tb4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf4 = tb4.text_frame
tf4.word_wrap = True

items4 = [
    ("Thách thức 1: Tranh chấp Mua vé Đồng thời (Race Condition)", 
     "Giải pháp: Áp dụng Pessimistic Concurrency Control (WITH (UPDLOCK, HOLDLOCK)) + IsolationLevel.Serializable + Cập nhật điều kiện nguyên tử (Atomic Update)."),
    ("Thách thức 2: Quản lý Đơn hàng Quá hạn & Khóa Giữ chỗ", 
     "Giải pháp: Cấu hình khóa giữ vé 10 phút, kích hoạt tự động tiến trình GiaiPhongDonHangHetHan hoàn trả kho vé/ghế khi hết hạn."),
    ("Thách thức 3: Bảo mật Thanh toán & Chống giả mạo Callback", 
     "Giải pháp: Xác thực chữ ký điện tử hai chiều HMAC-SHA256 / SHA256 Checksum, đảm bảo tính Toàn vẹn (Integrity) và Chống chối bỏ."),
    ("Thách thức 4: Cổng Soát vé Gian lận (Double Check-in)", 
     "Giải pháp: Mỗi vé phát hành 1 mã QR duy nhất mã hóa. Hệ thống chặn soát vé lần 2 và ghi vết lịch sử nhân viên check-in real-time.")
]

for idx, (title, desc) in enumerate(items4):
    p = tf4.paragraphs[0] if idx == 0 else tf4.add_paragraph()
    p.space_after = Pt(10)
    
    run1 = p.add_run()
    run1.text = f"• {title}\n"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(15)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_ACCENT
    
    run2 = p.add_run()
    run2.text = f"   ➔ {desc}"
    run2.font.name = FONT_NAME
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 5: NGHIỆP VỤ + THIẾT KẾ CSDL QUAN HỆ
# ==========================================
slide5 = prs.slides[4]
clear_shapes_text(slide5)
add_header(slide5, "4. NGHIỆP VỤ VÀ THIẾT KẾ DỮ LIỆU (CSDL QUAN HỆ)", "Mô hình Cơ sở Dữ liệu 11 bảng chuẩn hóa và các ràng buộc toàn vẹn")

tb5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf5 = tb5.text_frame
tf5.word_wrap = True

items5 = [
    ("Cấu trúc Mô hình 11 Bảng Quan hệ:", 
     "\n   - Nhóm Tài khoản & Phân quyền: NguoiDung, ThongTinBanToChuc, NhanVienSuKien."
     "\n   - Nhóm Sự kiện & Ghế ngồi: SuKien, LoaiVe, SoDoChoNgoi, KhuVuc, HangGhe, ChoNgoi."
     "\n   - Nhóm Giao dịch & Vé điện tử: DonHang, ChiTietDonHang (MaVe, MaQR, TrangThaiCheckin), MaGiamGia."),
    ("Các Ràng buộc Toàn vẹn Chính (Database Constraints):", 
     "\n   - Primary Key & Foreign Key (CASCADE / RESTRICT): Bảo đảm tính toàn vẹn tham chiếu."
     "\n   - Ràng buộc Duy nhất (UNIQUE): Đảm bảo duy nhất trên Email, MaDonHang, MaVe, MaQR, MaCode."
     "\n   - Ràng buộc Kiểm tra (CHECK): GiaBan >= 0, (SoLuongDaBan + SoLuongGiuCho) <= SoLuongTong.")
]

for idx, (title, desc) in enumerate(items5):
    p = tf5.paragraphs[0] if idx == 0 else tf5.add_paragraph()
    p.space_after = Pt(12)
    
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT_NAME
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 6: GIAO DIỆN HỆ THỐNG
# ==========================================
slide6 = prs.slides[5]
clear_shapes_text(slide6)
add_header(slide6, "5. GIAO DIỆN HỆ THỐNG (SYSTEM INTERFACES)", "Tổng quan các nhóm màn hình giao diện người dùng chính")

tb6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf6 = tb6.text_frame
tf6.word_wrap = True

items6 = [
    ("1. Nhóm Giao diện Khách hàng (Buyer UI):", " Trang chủ sự kiện nổi bật, Chi tiết sự kiện & Sơ đồ chọn ghế trực quan, Màn hình Thanh toán 3 cổng & Chi tiết Vé điện tử QR Code."),
    ("2. Nhóm Giao diện Ban tổ chức (Organizer UI):", " Dashboard thống kê doanh thu / tỷ lệ check-in, Giao diện Tạo sự kiện & Cấu hình loại vé/ghế, Quản lý Voucher & Yêu cầu rút tiền."),
    ("3. Nhóm Giao diện Nhân viên Soát vé (Staff UI):", " Màn hình Web Camera Quét mã QR Code & Nhập mã vé thủ công, Phản hồi màu Xanh (Thành công) / Đỏ (Vé đã check-in / Lỗi)."),
    ("4. Nhóm Giao diện Quản trị viên (Admin UI):", " Dashboard tổng quan toàn hệ thống, Màn hình Duyệt hồ sơ Ban tổ chức & Duyệt sự kiện mở bán, Quản lý đơn hàng & Hoàn tiền nội bộ.")
]

for idx, (title, desc) in enumerate(items6):
    p = tf6.paragraphs[0] if idx == 0 else tf6.add_paragraph()
    p.space_after = Pt(14)
    
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = f"\n   {desc}"
    run2.font.name = FONT_NAME
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 7: CÁC CHỨC NĂNG ĐÃ XÂY DỰNG
# ==========================================
slide7 = prs.slides[6]
clear_shapes_text(slide7)
add_header(slide7, "6. CÁC CHỨC NĂNG ĐÃ XÂY DỰNG", "Danh mục danh sách chức năng theo từng tác nhân (Chỉ liệt kê gạch đầu dòng)")

tb7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf7 = tb7.text_frame
tf7.word_wrap = True

items7 = [
    ("Khách hàng (Buyer):", " Tìm kiếm/lọc sự kiện, Đăng ký/Đăng nhập & Google OAuth, Đặt vé thường/ghế sơ đồ, Áp dụng Voucher, Thanh toán VNPay/MoMo/ZaloPay, Xem vé QR."),
    ("Ban tổ chức (Organizer):", " Đăng ký hồ sơ, Tạo/sửa sự kiện, Cấu hình loại vé & sơ đồ ghế, Quản lý Voucher, Phân công Staff, Báo cáo doanh thu/check-in, Xuất CSV, Rút tiền."),
    ("Nhân viên Soát vé (Staff):", " Truy cập sự kiện được phân công, Quét mã QR bằng Web Camera, Nhập mã vé thủ công, Kiểm tra trạng thái vé real-time."),
    ("Quản trị viên (Admin):", " Duyệt hồ sơ Ban tổ chức, Duyệt sự kiện mở bán, Quản lý đơn hàng toàn hệ thống, Hoàn tiền nội bộ, Duyệt rút tiền, Quản lý cấu hình & Audit Log.")
]

for idx, (title, desc) in enumerate(items7):
    p = tf7.paragraphs[0] if idx == 0 else tf7.add_paragraph()
    p.space_after = Pt(12)
    
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_ACCENT
    
    run2 = p.add_run()
    run2.text = f"\n   {desc}"
    run2.font.name = FONT_NAME
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 8: KẾT QUẢ THỰC HIỆN
# ==========================================
slide8 = prs.slides[7]
clear_shapes_text(slide8)
add_header(slide8, "7. KẾT QUẢ THỰC HIỆN & CHỈ SỐ NĂNG LỰC", "Đánh giá kết quả xây dựng dự án và các thông số chịu tải kỹ thuật")

tb8 = slide8.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf8 = tb8.text_frame
tf8.word_wrap = True

items8 = [
    ("Hoàn thành 100% Khung ứng dụng Web:", " Website chạy ổn định trên .NET 9.0, tương thích Responsive trên Desktop & Mobile."),
    ("Số người dùng truy cập đồng thời (Concurrent Users):", " Đạt từ 3,000 – 5,000 người dùng online cùng lúc nhờ Kestrel ThreadPool Async."),
    ("Tải truy vấn đọc (Read Throughput - RPS):", " Đạt ~2,500 Requests/giây nhờ Dapper SQL có Indexing tối ưu."),
    ("Tải giao dịch ghi (Write Throughput - TPS):", " Đạt ~150 – 300 Giao dịch/giây, đảm bảo hàng chờ an toàn tại SQL Server."),
    ("Chỉ số An toàn Dữ liệu Đồng thời:", " Đạt 0% Bán vượt kho (Overselling), 0% Trùng ghế (Double Booking) và 100% Bảo mật chữ ký thanh toán.")
]

for idx, (title, desc) in enumerate(items8):
    p = tf8.paragraphs[0] if idx == 0 else tf8.add_paragraph()
    p.space_after = Pt(12)
    
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT_NAME
    run2.font.size = Pt(15)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 9: KẾT LUẬN & HƯỚNG MỞ RỘNG
# ==========================================
slide9 = prs.slides[8]
clear_shapes_text(slide9)
add_header(slide9, "8. KẾT LUẬN VÀ HƯỚNG MỞ RỘNG", "Đánh giá mục tiêu đối sánh, hạn chế tồn tại và hướng phát triển tương lai")

# Create a small summary table or clear bullet points
tb9 = slide9.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.5), Inches(4.8))
tf9 = tb9.text_frame
tf9.word_wrap = True

items9 = [
    ("1. Đối sánh Mục tiêu vs Kết quả:", 
     " Đã hoàn thành toàn bộ mục tiêu đề ra (Thanh toán 3 cổng, Giữ chỗ 10 phút, Vé QR Code, Khóa bi quan chống overbooking)."),
    ("2. Những Hạn chế / Chưa làm được:", 
     " Chưa tích hợp API hoàn tiền tự động trực tiếp từ ngân hàng (mới ghi nhận nội bộ); Chưa có ứng dụng Mobile Native."),
    ("3. Hướng Phát triển Mở rộng:", 
     " Tích hợp Hàng chờ Redis / RabbitMQ để nâng công suất ghi lên 2,000+ TPS; Xây dựng Ứng dụng Mobile App (Flutter) quét QR Code.")
]

for idx, (title, desc) in enumerate(items9):
    p = tf9.paragraphs[0] if idx == 0 else tf9.add_paragraph()
    p.space_after = Pt(14)
    
    run1 = p.add_run()
    run1.text = f"• {title}"
    run1.font.name = FONT_NAME
    run1.font.size = Pt(16)
    run1.font.bold = True
    run1.font.color.rgb = COLOR_PRIMARY
    
    run2 = p.add_run()
    run2.text = desc
    run2.font.name = FONT_NAME
    run2.font.size = Pt(15)
    run2.font.color.rgb = COLOR_TEXT_MAIN

# ==========================================
# SLIDE 10: CẢM ƠN & Q&A
# ==========================================
slide10 = prs.slides[9]
clear_shapes_text(slide10)

tb10 = slide10.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(3.5))
tf10 = tb10.text_frame
tf10.word_wrap = True

p = tf10.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = "XIN CHÂN THÀNH CẢM ƠN QUÝ THẦY CÔ TRONG HỘI ĐỒNG ĐÃ LẮNG NGHE!"
p.font.name = FONT_NAME
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(30)

p = tf10.add_paragraph()
p.alignment = PP_ALIGN.CENTER
p.text = "PHIÊN THẢO LUẬN & HỎI ĐÁP (Q&A)"
p.font.name = FONT_NAME
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = COLOR_ACCENT

prs.save(output_path)
# Also update the template file
prs.save(template_path)

print(f"Successfully generated PowerPoint presentation at: {output_path}")
