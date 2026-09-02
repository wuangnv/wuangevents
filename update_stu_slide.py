import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

sys.stdout.reconfigure(encoding='utf-8')

filepath = 'C:/Users/Wuang/Desktop/Vibe/TRƯỜNG ĐẠI HỌC CÔNG NGHỆ SÀI GÒN KHOA CÔNG NGHỆ THÔNG TIN.pptx'
prs = Presentation(filepath)

# ----------------------------------------------------
# FIX & ENHANCE SLIDE 2: MỤC TIÊU ĐỀ TÀI
# ----------------------------------------------------
slide2 = prs.slides[1]
for shape in slide2.shapes:
    if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            if sub.has_text_frame:
                text = sub.text_frame.text.strip()
                if 'KHÁCH HÀNG' in text:
                    pass
                elif 'Tìm kiếm' in text:
                    sub.text_frame.text = "Số hóa trải nghiệm mua vé: Tìm kiếm, chọn vé thường/sơ đồ ghế, áp mã giảm giá, thanh toán 3 cổng & nhận vé QR."
                    sub.text_frame.paragraphs[0].font.size = Pt(13)
                elif 'BAN TỔ CHỨC' in text:
                    pass
                elif 'Quản lý sự kiện' in text:
                    sub.text_frame.text = "Tự động hóa quản lý vòng đời sự kiện: Tạo vé/ghế, quản lý ưu đãi, phân công staff, xem ví doanh thu & xuất CSV khách."
                    sub.text_frame.paragraphs[0].font.size = Pt(13)
                elif 'NHÂN VIÊN SOÁT VÉ' in text:
                    pass
                elif 'Quét QR' in text:
                    # FIX CUT-OFF TEXT HERE!
                    sub.text_frame.text = "Tối ưu hóa kiểm soát cửa vào: Quét mã QR bằng Web Camera hoặc nhập mã thủ công real-time, ngăn ngừa gian lận check-in."
                    sub.text_frame.paragraphs[0].font.size = Pt(13)
                elif 'QUẢN TRỊ VIÊN' in text:
                    pass
                elif 'Duyệt hồ sơ' in text:
                    sub.text_frame.text = "Kiểm soát an toàn toàn hệ thống: Duyệt hồ sơ Organizer & sự kiện mở bán, quản lý người dùng, giám sát đơn & hoàn tiền."
                    sub.text_frame.paragraphs[0].font.size = Pt(13)
                elif 'MỤC TIÊU KỸ THUẬT' in text:
                    pass
                elif 'Hạn chế bán vượt vé' in text:
                    sub.text_frame.text = "Bảo đảm an toàn giao dịch đồng thời (0% Overselling & 0% Double Booking với Khóa bi quan UPDLOCK) • Bảo mật HMAC-SHA256 • Responsive UI."
                    sub.text_frame.paragraphs[0].font.size = Pt(13)

# ----------------------------------------------------
# POPULATE SLIDE 8: MỘT SỐ GIAO DIỆN TIÊU BIỂU
# ----------------------------------------------------
slide8 = prs.slides[7]
tb8 = slide8.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf8 = tb8.text_frame
tf8.word_wrap = True

items8 = [
    "• Màn hình Khách hàng: Trang chủ sự kiện nổi bật, Chi tiết sự kiện & Sơ đồ chọn ghế trực quan, Giao diện Thanh toán 3 cổng (VNPay/MoMo/ZaloPay) & Vé điện tử QR Code.",
    "• Màn hình Ban tổ chức: Dashboard thống kê doanh thu / tỷ lệ check-in real-time, Giao diện tạo sự kiện & cấu hình loại vé/ghế, Quản lý Voucher & Rút tiền.",
    "• Màn hình Nhân viên Soát vé: Giao diện Web Camera Quét mã QR Code & Nhập mã vé thủ công, Phản hồi màu Xanh (Thành công) / Đỏ (Vé đã check-in / Lỗi).",
    "• Màn hình Quản trị viên: Dashboard tổng quan hệ thống, Duyệt hồ sơ Organizer & Sự kiện mở bán, Quản lý đơn hàng & Hoàn tiền nội bộ."
]
for idx, text in enumerate(items8):
    p = tf8.paragraphs[0] if idx == 0 else tf8.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(31, 41, 55)
    p.space_after = Pt(14)

# ----------------------------------------------------
# POPULATE SLIDE 9: PHÂN HỆ 01 - KHÁCH HÀNG (CHỨC NĂNG)
# ----------------------------------------------------
slide9 = prs.slides[8]
tb9 = slide9.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf9 = tb9.text_frame
tf9.word_wrap = True

items9 = [
    "• Tìm kiếm & Lọc sự kiện theo từ khóa, danh mục, thành phố và giá vé.",
    "• Đăng ký, đăng nhập tài khoản & Đăng nhập nhanh bằng Google OAuth 2.0.",
    "• Chọn loại vé thường hoặc Chọn vị trí ghế ngồi trực quan trên Sơ đồ sơ đồ ghế.",
    "• Áp dụng Mã giảm giá (Voucher) & Thanh toán trực tuyến qua 3 cổng VNPay, MoMo, ZaloPay.",
    "• Quản lý Lịch sử đơn hàng & Hiển thị Vé điện tử chứa mã QR Code duy nhất để check-in."
]
for idx, text in enumerate(items9):
    p = tf9.paragraphs[0] if idx == 0 else tf9.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(31, 41, 55)
    p.space_after = Pt(12)

# ----------------------------------------------------
# POPULATE SLIDE 10: PHÂN HỆ 02 - BAN TỔ CHỨC (CHỨC NĂNG)
# ----------------------------------------------------
slide10 = prs.slides[9]
tb10 = slide10.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf10 = tb10.text_frame
tf10.word_wrap = True

items10 = [
    "• Đăng ký hồ sơ Ban tổ chức (Organizer) & Quản lý thông tin tài khoản nhận tiền.",
    "• Tạo mới, lưu nháp, chỉnh sửa sự kiện, cấu hình loại vé & dựng Sơ đồ ghế ngồi.",
    "• Quản lý mã giảm giá (Voucher): giảm %, giới hạn số lượt sử dụng & hạn dùng.",
    "• Phân công Nhân viên soát vé (Staff) theo email truy cập.",
    "• Xem Dashboard báo cáo doanh thu, tiến độ check-in, xuất file CSV khách tham dự & Gửi yêu cầu rút tiền."
]
for idx, text in enumerate(items10):
    p = tf10.paragraphs[0] if idx == 0 else tf10.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(31, 41, 55)
    p.space_after = Pt(12)

# ----------------------------------------------------
# POPULATE SLIDE 11: PHÂN HỆ 03-04 - VẬN HÀNH (STAFF & ADMIN)
# ----------------------------------------------------
slide11 = prs.slides[10]
tb11 = slide11.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf11 = tb11.text_frame
tf11.word_wrap = True

items11 = [
    "• PHÂN HỆ SOÁT VÉ (STAFF):",
    "   - Mở giao diện check-in sự kiện được phân công.",
    "   - Quét mã QR Code qua Web Camera hoặc Nhập mã vé thủ công.",
    "   - Kiểm tra trạng thái vé real-time, ngăn chặn việc check-in lặp lọt 2 lần.",
    "• PHÂN HỆ QUẢN TRỊ VIÊN (ADMIN):",
    "   - Xét duyệt hồ sơ Organizer & Duyệt sự kiện mở bán công khai.",
    "   - Quản lý đơn hàng toàn hệ thống & Thực hiện hoàn tiền nội bộ.",
    "   - Duyệt yêu cầu rút tiền của Organizer, Quản lý cấu hình & Nhật ký thao tác (Audit Log)."
]
for idx, text in enumerate(items11):
    p = tf11.paragraphs[0] if idx == 0 else tf11.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(15 if text.startswith('   -') else 17)
    p.font.bold = not text.startswith('   -')
    p.font.color.rgb = RGBColor(37, 99, 235) if not text.startswith('   -') else RGBColor(31, 41, 55)
    p.space_after = Pt(8)

# ----------------------------------------------------
# POPULATE SLIDE 12: KẾT QUẢ THỰC HIỆN
# ----------------------------------------------------
slide12 = prs.slides[11]
tb12 = slide12.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf12 = tb12.text_frame
tf12.word_wrap = True

items12 = [
    "• Hoàn thành 100% Khung ứng dụng Web responsive trên nền .NET 9.0 ASP.NET Core MVC.",
    "• Đạt năng lực phục vụ 3,000 – 5,000 người dùng truy cập đồng thời.",
    "• Tải truy vấn đọc (Read Throughput): ~2,500 Requests/giây nhờ Dapper SQL có Indexing.",
    "• Tải giao dịch ghi (Write Throughput): ~150 – 300 Giao dịch/giây với SQL Lock Queue.",
    "• Bảo đảm an toàn tuyệt đối 0% Overselling (không bán vượt kho) & 0% Double Booking (không trùng ghế)."
]
for idx, text in enumerate(items12):
    p = tf12.paragraphs[0] if idx == 0 else tf12.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(31, 41, 55)
    p.space_after = Pt(12)

# ----------------------------------------------------
# POPULATE SLIDE 14: HẠN CHẾ VÀ HƯỚNG PHÁT TRIỂN
# ----------------------------------------------------
slide14 = prs.slides[13]
tb14 = slide14.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(16.0), Inches(7.0))
tf14 = tb14.text_frame
tf14.word_wrap = True

items14 = [
    "• NHỮNG HẠN CHẾ CÒN TỒN TẠI:",
    "   - Chưa tích hợp API hoàn tiền tự động trực tiếp từ cổng ngân hàng (mới ghi nhận nội bộ).",
    "   - Chưa xây dựng ứng dụng di động Native (Mobile App) riêng.",
    "• HƯỚNG PHÁT TRIỂN NGHỆ MỞ RỘNG:",
    "   - Tích hợp Hàng chờ Redis / Message Queue (RabbitMQ) để nâng công suất ghi lên 2,000+ TPS.",
    "   - Xây dựng Ứng dụng Mobile App (Flutter / React Native) quét mã QR chuyên dụng.",
    "   - Tích hợp VietQR API để tự động hóa đối soát chuyển khoản ngân hàng."
]
for idx, text in enumerate(items14):
    p = tf14.paragraphs[0] if idx == 0 else tf14.add_paragraph()
    p.text = text
    p.font.name = 'Segoe UI'
    p.font.size = Pt(15 if text.startswith('   -') else 17)
    p.font.bold = not text.startswith('   -')
    p.font.color.rgb = RGBColor(37, 99, 235) if not text.startswith('   -') else RGBColor(31, 41, 55)
    p.space_after = Pt(8)

prs.save(filepath)
print("Successfully updated TRƯỜNG ĐẠI HỌC CÔNG NGHỆ SÀI GÒN KHOA CÔNG NGHỆ THÔNG TIN.pptx!")
